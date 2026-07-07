"""Beautify 7月生産検討会資料.pptx with consistent corporate styling."""
import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parents[1]
SRC = BASE / "7月生産検討会資料.pptx"
BACKUP = BASE / "7月生産検討会資料_原版.pptx"
OUT = BASE / "7月生産検討会資料.pptx"

# Corporate palette
PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_LIGHT = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2D, 0x34, 0x40)
TEXT_MUTED = RGBColor(0x5A, 0x6A, 0x7A)
TABLE_ALT = RGBColor(0xF4, 0xF7, 0xFB)
TABLE_BORDER = RGBColor(0xC5, 0xD3, 0xE0)
SUCCESS = RGBColor(0x1A, 0x7F, 0x4E)
WARNING = RGBColor(0xC9, 0x7A, 0x00)

FONT = "Meiryo UI"
PART_SLIDES = {2, 6, 9, 11, 14}
COVER_SLIDE = 1
SECTION_COVER = 15


def set_run_font(run, size=None, bold=None, color=None, name=FONT):
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def style_paragraph(p, size=14, bold=False, color=TEXT_DARK, align=None, space_after=6):
    p.alignment = align if align is not None else p.alignment
    p.space_after = Pt(space_after)
    for run in p.runs:
        set_run_font(run, size=size, bold=bold, color=color)


def style_text_frame(tf, default_size=14, default_color=TEXT_DARK, bold=False, align=None):
    if not tf:
        return
    tf.word_wrap = True
    for p in tf.paragraphs:
        style_paragraph(p, size=default_size, bold=bold, color=default_color, align=align)


def is_title_shape(shape, slide_idx):
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if slide_idx in PART_SLIDES and text.startswith("PART "):
        return True
    if len(text) > 60:
        return False
    # left accent bar companion titles
    if shape.top and shape.top.inches < 1.2 and len(text) < 40:
        return True
    return False


def is_subtitle_shape(shape, slide_idx):
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if slide_idx in PART_SLIDES and len(text) > 20 and not text.startswith("PART"):
        return True
    if slide_idx == COVER_SLIDE and "各工程" in text:
        return True
    if slide_idx == SECTION_COVER and text.startswith("〜"):
        return True
    return False


def style_table(table):
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = cell.fill
            if r_idx == 0 or (r_idx == 1 and c_idx == 0 and cell.text.strip() == ""):
                fill.solid()
                fill.fore_color.rgb = PRIMARY if r_idx == 0 else ACCENT_LIGHT
            elif r_idx % 2 == 0:
                fill.solid()
                fill.fore_color.rgb = TABLE_ALT
            else:
                fill.solid()
                fill.fore_color.rgb = WHITE

            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r_idx <= 1 else PP_ALIGN.CENTER
                size = 10 if r_idx <= 1 else 11
                bold = r_idx <= 1
                color = WHITE if r_idx == 0 else TEXT_DARK
                for run in p.runs:
                    set_run_font(run, size=size, bold=bold, color=color)


def clean_text(text: str) -> str:
    text = re.sub(r"[\u3000\s]{3,}", " ", text)
    text = text.replace("\u27a1\ufe0f", "→").replace("\u27a1", "→")
    return text.strip()


def remove_bullet_shapes(slide):
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "■":
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def style_accent_bar(shape):
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return
    if shape.width.inches < 0.2 and shape.height.inches < 1.0:
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()


def style_part_slide(slide, slide_idx):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = clean_text(shape.text_frame.text)
        if text.startswith("PART "):
            tf = shape.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in tf.paragraphs[0].runs:
                set_run_font(run, size=32, bold=True, color=PRIMARY)
        elif len(text) > 15 and shape.top.inches > 3.5:
            tf = shape.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in tf.paragraphs[0].runs:
                set_run_font(run, size=16, bold=False, color=TEXT_MUTED)


def style_cover_slide(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = clean_text(shape.text_frame.text)
        if "生産検討会" in text:
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=40, bold=True, color=PRIMARY)
        elif "各工程" in text:
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=18, bold=False, color=TEXT_MUTED)
        elif re.match(r"\d{4}年", text) or "生産管理部" in text:
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=14, bold=False, color=TEXT_MUTED)


def style_content_title(shape):
    text = clean_text(shape.text_frame.text)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    for run in shape.text_frame.paragraphs[0].runs:
        set_run_font(run, size=24, bold=True, color=PRIMARY)


def style_summary_box(shape):
    """Bottom summary text boxes."""
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.space_before = Pt(2)
        for run in p.runs:
            set_run_font(run, size=12, bold=False, color=TEXT_DARK)


def style_analysis_cards(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = clean_text(shape.text_frame.text)
        if text.startswith("分析："):
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=10, bold=False, color=TEXT_MUTED)
        elif text.startswith("対策："):
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=10, bold=True, color=ACCENT)
        elif "負荷率" in text and text.endswith("）"):
            for run in shape.text_frame.paragraphs[0].runs:
                set_run_font(run, size=11, bold=True, color=PRIMARY)


def add_footer(slide, slide_idx, total):
    if slide_idx in PART_SLIDES or slide_idx == COVER_SLIDE or slide_idx == SECTION_COVER:
        return
    left = Inches(0.5)
    top = Inches(7.15)
    width = Inches(12.3)
    height = Inches(0.25)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"7月度 生産検討会  |  {slide_idx} / {total}"
    set_run_font(run, size=9, color=TEXT_MUTED)


def beautify():
    if not BACKUP.exists():
        shutil.copy2(SRC, BACKUP)
        print(f"Backup saved: {BACKUP.name}")

    prs = Presentation(str(SRC))
    total = len(prs.slides)

    for i, slide in enumerate(prs.slides, 1):
        remove_bullet_shapes(slide)

        for shape in slide.shapes:
            style_accent_bar(shape)

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                style_table(shape.table)
                continue

            if not shape.has_text_frame:
                continue

            # clean whitespace
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = clean_text(run.text)

            if i == COVER_SLIDE:
                continue
            if i in PART_SLIDES:
                continue
            if i == SECTION_COVER:
                continue

            text = clean_text(shape.text_frame.text)
            if not text:
                continue

            if shape.top and shape.top.inches < 1.0 and len(text) < 45:
                style_content_title(shape)
            elif text.startswith(("全体動向", "工程別", "特記事項", "在庫への影響", "改善効果", "コスト効果")):
                style_summary_box(shape)
            elif text.startswith(("出荷数", "生産性向上")):
                style_summary_box(shape)
            elif "稼働日" in text and "内示" in text:
                for run in shape.text_frame.paragraphs[0].runs:
                    set_run_font(run, size=11, bold=True, color=ACCENT)

        if i == COVER_SLIDE:
            style_cover_slide(slide)
        elif i in PART_SLIDES:
            style_part_slide(slide, i)
        elif i in (7, 10):
            style_analysis_cards(slide)

        add_footer(slide, i, total)

    prs.save(str(OUT))
    print(f"Beautified: {OUT.name}")


if __name__ == "__main__":
    beautify()
