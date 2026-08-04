"""実績修正統計 PPT 生成。

画面（ManualEntryStatistics）の KPI・月次比較・推移・工程別データを
ワイド画面スライドに出力する。グラフは Pillow で描画し埋め込む。
"""
from __future__ import annotations

import io
import math
from datetime import datetime
from typing import Any, Dict, List, Optional, Sequence, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x0F, 0x2B, 0x4C)
ACCENT = RGBColor(0x1D, 0x6F, 0xC2)
ACCENT_SOFT = RGBColor(0xD8, 0xE8, 0xF8)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ORANGE = RGBColor(0xE0, 0x7A, 0x1A)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
ROSE = RGBColor(0xE1, 0x1D, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TABLE_ALT = RGBColor(0xF1, 0xF5, 0xF9)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
FONT = "Meiryo UI"

# チャート配色（画面に近い）
C_PROD = (37, 99, 235)  # 実績修正
C_AUTO = (13, 148, 136)  # 実績集計
C_CMP = (148, 163, 184)  # 比較月
C_RATIO = (234, 88, 12)  # 比率
C_GRID = (226, 232, 240)
C_AXIS = (100, 116, 139)
C_BG = (255, 255, 255)
C_PANEL = (248, 250, 252)


def _fmt_num(v: Any) -> str:
    try:
        return f"{int(round(float(v))):,}"
    except (TypeError, ValueError):
        return "0"


def _fmt_qty(v: Any) -> str:
    try:
        n = float(v)
    except (TypeError, ValueError):
        return "0"
    if abs(n - round(n)) < 1e-9:
        return f"{int(round(n)):,}"
    return f"{n:,.1f}"


def _fmt_qty_sen(v: Any) -> str:
    try:
        return f"{float(v) / 1000:,.1f}"
    except (TypeError, ValueError):
        return "0.0"


def _fmt_pct(v: Any) -> str:
    try:
        return f"{float(v) * 100:.1f}%"
    except (TypeError, ValueError):
        return "0.0%"


def _fmt_delta(change: Any, rate: Any = None) -> str:
    try:
        c = float(change)
    except (TypeError, ValueError):
        return "—"
    sign = "+" if c > 0 else ""
    base = f"{sign}{_fmt_num(c)}"
    if rate is None:
        return base
    try:
        r = float(rate) * 100
        return f"{base} ({'+' if r > 0 else ''}{r:.1f}%)"
    except (TypeError, ValueError):
        return base


def _fmt_pct_point(v: Any) -> str:
    try:
        n = float(v) * 100
    except (TypeError, ValueError):
        return "—"
    return f"{'+' if n >= 0 else ''}{n:.1f}pt"


def _font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/meiryob.ttc" if bold else "C:/Windows/Fonts/meiryo.ttc",
        "C:/Windows/Fonts/YuGothB.ttc" if bold else "C:/Windows/Fonts/YuGothM.ttc",
        "C:/Windows/Fonts/msgothic.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size=size, index=0)
        except OSError:
            continue
    return ImageFont.load_default()


def _text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> Tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def _rounded_rect(draw: ImageDraw.ImageDraw, box, radius: int, fill) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill)


def _nice_max(values: Sequence[float]) -> float:
    positive = [float(v) for v in values if v is not None and float(v) > 0]
    if not positive:
        return 1.0
    m = max(positive)
    if m <= 1:
        return 1.0
    exp = 10 ** math.floor(math.log10(m))
    for step in (1, 1.25, 1.5, 2, 2.5, 5, 10):
        cand = step * exp
        if cand >= m:
            return cand
    return math.ceil(m / exp) * exp


def _draw_legend(draw: ImageDraw.ImageDraw, items: List[Tuple[Tuple[int, int, int], str]], x: int, y: int) -> None:
    font = _font(22, bold=True)
    cx = x
    for color, label in items:
        draw.rounded_rectangle((cx, y + 4, cx + 22, y + 22), radius=4, fill=color)
        tw, _ = _text_size(draw, label, font)
        draw.text((cx + 30, y), label, font=font, fill=(30, 41, 59))
        cx += 30 + tw + 28


def _draw_grouped_bars(
    categories: Sequence[str],
    series: List[Tuple[str, Sequence[float], Tuple[int, int, int]]],
    *,
    width: int = 1400,
    height: int = 720,
    value_fmt=None,
    title: str = "",
) -> bytes:
    img = Image.new("RGB", (width, height), C_BG)
    draw = ImageDraw.Draw(img)
    _rounded_rect(draw, (0, 0, width - 1, height - 1), 18, C_PANEL)
    _rounded_rect(draw, (2, 2, width - 3, height - 3), 16, C_BG)

    pad_l, pad_r, pad_t, pad_b = 90, 40, 70, 90
    if title:
        draw.text((28, 18), title, font=_font(28, bold=True), fill=(15, 43, 76))
        pad_t = 88

    legend_items = [(c, n) for n, _, c in series]
    _draw_legend(draw, legend_items, pad_l, pad_t - 42)

    plot_l, plot_t = pad_l, pad_t
    plot_r, plot_b = width - pad_r, height - pad_b
    plot_w = plot_r - plot_l
    plot_h = plot_b - plot_t

    all_vals = [float(v) for _, vals, _ in series for v in vals]
    vmax = _nice_max(all_vals)
    axis_font = _font(18)
    label_font = _font(20, bold=True)
    value_font = _font(16, bold=True)

    for i in range(5):
        y = plot_b - int(plot_h * i / 4)
        val = vmax * i / 4
        draw.line((plot_l, y, plot_r, y), fill=C_GRID, width=2)
        label = value_fmt(val) if value_fmt else _fmt_num(val)
        tw, th = _text_size(draw, label, axis_font)
        draw.text((plot_l - tw - 12, y - th // 2), label, font=axis_font, fill=C_AXIS)

    n = max(len(categories), 1)
    group_w = plot_w / n
    bar_n = max(len(series), 1)
    bar_w = max(14, int(group_w * 0.62 / bar_n))
    gap = 6

    for gi, cat in enumerate(categories):
        gx = plot_l + group_w * gi + group_w / 2
        total_bars_w = bar_n * bar_w + (bar_n - 1) * gap
        start_x = gx - total_bars_w / 2
        for si, (_, vals, color) in enumerate(series):
            v = float(vals[gi]) if gi < len(vals) else 0.0
            bh = int((v / vmax) * (plot_h - 8)) if vmax else 0
            x0 = int(start_x + si * (bar_w + gap))
            y0 = plot_b - bh
            # soft top highlight
            draw.rounded_rectangle((x0, y0, x0 + bar_w, plot_b), radius=6, fill=color)
            if bh > 18:
                hi = tuple(min(255, c + 40) for c in color)
                draw.rounded_rectangle((x0 + 2, y0 + 2, x0 + bar_w - 2, y0 + 10), radius=4, fill=hi)
            if v > 0:
                txt = value_fmt(v) if value_fmt else _fmt_num(v)
                tw, th = _text_size(draw, txt, value_font)
                draw.text((x0 + bar_w / 2 - tw / 2, y0 - th - 4), txt, font=value_font, fill=color)
        tw, th = _text_size(draw, str(cat), label_font)
        draw.text((gx - tw / 2, plot_b + 14), str(cat), font=label_font, fill=(30, 41, 59))

    draw.line((plot_l, plot_b, plot_r, plot_b), fill=(148, 163, 184), width=3)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _draw_combo_trend(
    categories: Sequence[str],
    bars: List[Tuple[str, Sequence[float], Tuple[int, int, int]]],
    line: Tuple[str, Sequence[float], Tuple[int, int, int]],
    *,
    width: int = 1400,
    height: int = 720,
    line_as_pct: bool = True,
    title: str = "",
) -> bytes:
    img = Image.new("RGB", (width, height), C_BG)
    draw = ImageDraw.Draw(img)
    _rounded_rect(draw, (0, 0, width - 1, height - 1), 18, C_PANEL)
    _rounded_rect(draw, (2, 2, width - 3, height - 3), 16, C_BG)

    pad_l, pad_r, pad_t, pad_b = 90, 90, 70, 90
    if title:
        draw.text((28, 18), title, font=_font(28, bold=True), fill=(15, 43, 76))
        pad_t = 88

    legend = [(c, n) for n, _, c in bars] + [(line[2], line[0])]
    _draw_legend(draw, legend, pad_l, pad_t - 42)

    plot_l, plot_t = pad_l, pad_t
    plot_r, plot_b = width - pad_r, height - pad_b
    plot_w = plot_r - plot_l
    plot_h = plot_b - plot_t

    bar_vals = [float(v) for _, vals, _ in bars for v in vals]
    vmax = _nice_max(bar_vals)
    line_vals = [float(v) for v in line[1]]
    lmax = max(line_vals) if line_vals else 0.0
    lmax = max(lmax * 1.25, 0.05) if line_as_pct else _nice_max(line_vals)

    axis_font = _font(18)
    label_font = _font(20, bold=True)

    for i in range(5):
        y = plot_b - int(plot_h * i / 4)
        draw.line((plot_l, y, plot_r, y), fill=C_GRID, width=2)
        left_label = _fmt_num(vmax * i / 4)
        tw, th = _text_size(draw, left_label, axis_font)
        draw.text((plot_l - tw - 10, y - th // 2), left_label, font=axis_font, fill=C_AXIS)
        if line_as_pct:
            right_label = f"{(lmax * i / 4) * 100:.0f}%"
        else:
            right_label = _fmt_num(lmax * i / 4)
        draw.text((plot_r + 10, y - th // 2), right_label, font=axis_font, fill=line[2])

    n = max(len(categories), 1)
    group_w = plot_w / n
    bar_n = max(len(bars), 1)
    bar_w = max(12, int(group_w * 0.5 / bar_n))
    gap = 5
    centers: List[float] = []

    for gi, cat in enumerate(categories):
        gx = plot_l + group_w * gi + group_w / 2
        centers.append(gx)
        total_bars_w = bar_n * bar_w + (bar_n - 1) * gap
        start_x = gx - total_bars_w / 2
        for si, (_, vals, color) in enumerate(bars):
            v = float(vals[gi]) if gi < len(vals) else 0.0
            bh = int((v / vmax) * (plot_h - 8)) if vmax else 0
            x0 = int(start_x + si * (bar_w + gap))
            y0 = plot_b - bh
            draw.rounded_rectangle((x0, y0, x0 + bar_w, plot_b), radius=5, fill=color)
        tw, th = _text_size(draw, str(cat)[-5:], label_font)
        draw.text((gx - tw / 2, plot_b + 14), str(cat)[-5:], font=label_font, fill=(30, 41, 59))

    points = []
    for gi, v in enumerate(line_vals):
        x = centers[gi] if gi < len(centers) else plot_l
        y = plot_b - int((float(v) / lmax) * (plot_h - 8)) if lmax else plot_b
        points.append((x, y))
    if len(points) >= 2:
        draw.line(points, fill=line[2], width=5)
    for (x, y), v in zip(points, line_vals):
        draw.ellipse((x - 7, y - 7, x + 7, y + 7), fill=line[2], outline=(255, 255, 255), width=3)
        txt = f"{float(v) * 100:.1f}%" if line_as_pct else _fmt_num(v)
        tw, th = _text_size(draw, txt, _font(16, bold=True))
        draw.text((x - tw / 2, y - th - 10), txt, font=_font(16, bold=True), fill=line[2])

    draw.line((plot_l, plot_b, plot_r, plot_b), fill=(148, 163, 184), width=3)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def _set_run(run, *, size: int, bold: bool = False, color: RGBColor = TEXT_DARK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 16,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=font_size, bold=bold, color=color)
    return box


def _add_rect(slide, left, top, width, height, fill: RGBColor, *, line: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.12
    except Exception:
        pass
    return shape


def _add_accent_bar(slide, color: RGBColor = ACCENT) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_picture(slide, png: bytes, left, top, width, height) -> None:
    slide.shapes.add_picture(io.BytesIO(png), left, top, width=width, height=height)


def _style_table(table, header_color: RGBColor = PRIMARY, body_size: int = 13) -> None:
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    if r_idx == 0:
                        run.font.size = Pt(max(body_size - 1, 12))
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.size = Pt(body_size)
                        run.font.bold = False
                        run.font.color.rgb = TEXT_DARK
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                fill.fore_color.rgb = TABLE_ALT
            else:
                fill.fore_color.rgb = WHITE


def _add_table(
    slide,
    rows: List[List[str]],
    left,
    top,
    width,
    height,
    *,
    header_color: RGBColor = PRIMARY,
    body_size: int = 13,
    col_widths: Optional[Sequence[float]] = None,
):
    if not rows or not rows[0]:
        return None
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    if col_widths and len(col_widths) == len(rows[0]):
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(int(width * (w / total)))
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            if c_idx == 0:
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
    _style_table(table, header_color=header_color, body_size=body_size)
    for r_idx, row in enumerate(rows):
        cell = table.cell(r_idx, 0)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = FONT
                if r_idx == 0:
                    run.font.size = Pt(max(body_size - 1, 12))
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = TEXT_DARK
    return table


def _build_cover(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.9), SLIDE_W, Inches(3.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY
    banner.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), SLIDE_W, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    month = data.get("month") or ""
    cmp = data.get("compareMonth") or ""
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(2.25),
        Inches(11.7),
        Inches(0.9),
        "実績修正統計",
        font_size=44,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(3.2),
        Inches(11.7),
        Inches(0.55),
        "実績修正 vs 実績集計の月次比較（手入力除外）",
        font_size=20,
        color=ACCENT_SOFT,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(3.9),
        Inches(11.7),
        Inches(0.45),
        f"対象月 {month}　／　比較月 {cmp}　／　推移 {data.get('trendMonths', 6)}ヶ月",
        font_size=18,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11.7),
        Inches(0.4),
        f"生成日時 {datetime.now().strftime('%Y-%m-%d %H:%M')}　｜　Smart-EMAP",
        font_size=14,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def _build_kpi_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ACCENT)
    cur = data.get("current") or {}
    cmp = data.get("compare") or {}
    mom = data.get("monthOverMonth") or {}
    month = data.get("month") or ""
    cmp_month = data.get("compareMonth") or ""

    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.22),
        Inches(12.5),
        Inches(0.5),
        f"KPIサマリ　{month} vs {cmp_month}",
        font_size=28,
        bold=True,
        color=PRIMARY,
    )

    cards = [
        (
            "実績修正",
            ACCENT,
            f"{_fmt_num(cur.get('prodDataMgmt', {}).get('count'))} 件",
            f"数量 {_fmt_qty(cur.get('prodDataMgmt', {}).get('quantity'))}",
            f"前月比 {_fmt_delta(mom.get('prodDataMgmtCountChange'), mom.get('prodDataMgmtCountChangeRate'))}",
        ),
        (
            "実績集計",
            TEAL,
            f"{_fmt_num(cur.get('auto', {}).get('count'))} 件",
            f"数量 {_fmt_qty(cur.get('auto', {}).get('quantity'))}",
            f"総件数 {_fmt_num(cur.get('total', {}).get('count'))}",
        ),
        (
            "修正比率",
            ORANGE,
            _fmt_pct(cur.get("prodDataMgmtCountRatio")),
            f"数量比率 {_fmt_pct(cur.get('prodDataMgmtQuantityRatio'))}",
            f"前月比 {_fmt_pct_point(mom.get('prodDataMgmtCountRatioChange'))}",
        ),
    ]

    x = 0.4
    for title, color, main, sub, foot in cards:
        _add_rect(slide, Inches(x), Inches(0.9), Inches(4.0), Inches(2.35), CARD_BG)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(0.9), Inches(0.12), Inches(2.35)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        _add_textbox(slide, Inches(x + 0.28), Inches(1.05), Inches(3.5), Inches(0.4), title, font_size=18, bold=True, color=color)
        _add_textbox(slide, Inches(x + 0.28), Inches(1.5), Inches(3.5), Inches(0.6), main, font_size=30, bold=True, color=TEXT_DARK)
        _add_textbox(slide, Inches(x + 0.28), Inches(2.15), Inches(3.5), Inches(0.35), sub, font_size=16, color=TEXT_MUTED)
        _add_textbox(slide, Inches(x + 0.28), Inches(2.55), Inches(3.5), Inches(0.35), foot, font_size=15, bold=True, color=TEXT_DARK)
        x += 4.2

    # VS panel
    _add_rect(slide, Inches(0.4), Inches(3.55), Inches(6.15), Inches(3.5), CARD_BG)
    _add_rect(slide, Inches(6.8), Inches(3.55), Inches(6.15), Inches(3.5), CARD_BG)
    _add_textbox(slide, Inches(0.65), Inches(3.7), Inches(5.6), Inches(0.4), f"対象月 {month}", font_size=20, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(7.05), Inches(3.7), Inches(5.6), Inches(0.4), f"比較月 {cmp_month}", font_size=20, bold=True, color=TEXT_MUTED)

    def _vs_block(left: float, summary: dict) -> None:
        rows = [
            ("実績修正（件）", _fmt_num(summary.get("prodDataMgmt", {}).get("count"))),
            ("実績集計（件）", _fmt_num(summary.get("auto", {}).get("count"))),
            ("修正比率", _fmt_pct(summary.get("prodDataMgmtCountRatio"))),
            ("実績修正数量", _fmt_qty(summary.get("prodDataMgmt", {}).get("quantity"))),
            ("実績集計数量", _fmt_qty(summary.get("auto", {}).get("quantity"))),
            ("修正数量比率", _fmt_pct(summary.get("prodDataMgmtQuantityRatio"))),
        ]
        y = 4.25
        for label, val in rows:
            _add_textbox(slide, Inches(left), Inches(y), Inches(3.0), Inches(0.35), label, font_size=15, color=TEXT_MUTED)
            _add_textbox(slide, Inches(left + 2.7), Inches(y), Inches(2.8), Inches(0.35), val, font_size=18, bold=True, color=TEXT_DARK)
            y += 0.42

    _vs_block(0.65, cur)
    _vs_block(7.05, cmp)


def _build_month_compare_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ACCENT)
    month = data.get("month") or ""
    cmp_month = data.get("compareMonth") or ""
    cur = data.get("current") or {}
    cmp = data.get("compare") or {}
    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.2),
        Inches(12.5),
        Inches(0.45),
        "月次比較（件数・数量）",
        font_size=28,
        bold=True,
        color=PRIMARY,
    )

    cats = ["実績修正", "実績集計"]
    count_png = _draw_grouped_bars(
        cats,
        [
            ("対象月", [cur.get("prodDataMgmt", {}).get("count", 0), cur.get("auto", {}).get("count", 0)], C_PROD),
            ("比較月", [cmp.get("prodDataMgmt", {}).get("count", 0), cmp.get("auto", {}).get("count", 0)], C_CMP),
        ],
        width=1280,
        height=680,
        title=f"件数比較  {month} vs {cmp_month}",
    )
    qty_png = _draw_grouped_bars(
        cats,
        [
            (
                "対象月",
                [
                    float(cur.get("prodDataMgmt", {}).get("quantity", 0)) / 1000,
                    float(cur.get("auto", {}).get("quantity", 0)) / 1000,
                ],
                C_PROD,
            ),
            (
                "比較月",
                [
                    float(cmp.get("prodDataMgmt", {}).get("quantity", 0)) / 1000,
                    float(cmp.get("auto", {}).get("quantity", 0)) / 1000,
                ],
                C_CMP,
            ),
        ],
        width=1280,
        height=680,
        value_fmt=lambda v: f"{v:,.1f}",
        title=f"数量比較（千）  {month} vs {cmp_month}",
    )
    _add_picture(slide, count_png, Inches(0.35), Inches(0.8), Inches(6.3), Inches(6.3))
    _add_picture(slide, qty_png, Inches(6.75), Inches(0.8), Inches(6.3), Inches(6.3))


def _build_trend_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ORANGE)
    trend = data.get("byMonthTrend") or []
    months = [str(r.get("month") or "") for r in trend]
    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.2),
        Inches(12.5),
        Inches(0.45),
        f"修正比率推移（直近 {data.get('trendMonths', 6)} ヶ月）",
        font_size=28,
        bold=True,
        color=PRIMARY,
    )

    count_png = _draw_combo_trend(
        months,
        [
            ("実績修正", [r.get("prodDataMgmtCount", 0) for r in trend], C_PROD),
            ("実績集計", [r.get("autoCount", 0) for r in trend], C_AUTO),
        ],
        ("修正比率", [r.get("prodDataMgmtCountRatio", 0) for r in trend], C_RATIO),
        width=1280,
        height=680,
        title="件数 × 修正比率",
    )
    qty_png = _draw_combo_trend(
        months,
        [
            ("修正数量(千)", [float(r.get("prodDataMgmtQuantity", 0)) / 1000 for r in trend], C_PROD),
            ("集計数量(千)", [float(r.get("autoQuantity", 0)) / 1000 for r in trend], C_AUTO),
        ],
        ("数量比率", [r.get("prodDataMgmtQuantityRatio", 0) for r in trend], C_RATIO),
        width=1280,
        height=680,
        title="数量（千）× 修正数量比率",
    )
    _add_picture(slide, count_png, Inches(0.35), Inches(0.8), Inches(6.3), Inches(6.3))
    _add_picture(slide, qty_png, Inches(6.75), Inches(0.8), Inches(6.3), Inches(6.3))


def _build_process_chart_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, TEAL)
    month = data.get("month") or ""
    cmp_month = data.get("compareMonth") or ""
    cur_list = data.get("byProcess") or []
    cmp_map = {
        (r.get("processCd") or ""): r for r in (data.get("byProcessCompare") or [])
    }

    # merge top processes by current prod count
    rows = cur_list[:12]
    names = [str(r.get("processName") or r.get("processCd") or "")[:8] for r in rows]
    cur_prod = [r.get("prodDataMgmt", {}).get("count", 0) for r in rows]
    cmp_prod = [
        (cmp_map.get(r.get("processCd") or {}) or {}).get("prodDataMgmt", {}).get("count", 0) for r in rows
    ]
    cur_auto = [r.get("auto", {}).get("count", 0) for r in rows]

    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.2),
        Inches(12.5),
        Inches(0.45),
        f"工程別比較・内訳（上位 {len(rows)} 工程）",
        font_size=28,
        bold=True,
        color=PRIMARY,
    )

    compare_png = _draw_grouped_bars(
        names,
        [
            (f"修正 {month}", cur_prod, C_PROD),
            (f"修正 {cmp_month}", cmp_prod, C_CMP),
        ],
        width=1600,
        height=620,
        title=f"工程別 実績修正件数  {cmp_month} vs {month}",
    )
    breakdown_png = _draw_grouped_bars(
        names,
        [
            ("実績修正", cur_prod, C_PROD),
            ("実績集計", cur_auto, C_AUTO),
        ],
        width=1600,
        height=620,
        title=f"工程別内訳（件数）  {month}",
    )
    _add_picture(slide, compare_png, Inches(0.35), Inches(0.75), Inches(12.6), Inches(3.2))
    _add_picture(slide, breakdown_png, Inches(0.35), Inches(4.05), Inches(12.6), Inches(3.2))


def _build_process_table_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, VIOLET)
    month = data.get("month") or ""
    cmp_month = data.get("compareMonth") or ""
    cur_list = data.get("byProcess") or []
    cmp_map = {(r.get("processCd") or ""): r for r in (data.get("byProcessCompare") or [])}

    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.2),
        Inches(12.5),
        Inches(0.45),
        f"工程別比較一覧　{cmp_month} → {month}",
        font_size=26,
        bold=True,
        color=PRIMARY,
    )

    header = [
        "工程",
        f"修正件数\n{month}",
        f"修正件数\n{cmp_month}",
        "件数差",
        f"修正数量\n{month}",
        "比率",
        f"集計件数\n{month}",
    ]
    body: List[List[str]] = [header]
    keys = list({*(r.get("processCd") or "" for r in cur_list), *cmp_map.keys()})
    # order by current prod count
    cur_map = {(r.get("processCd") or ""): r for r in cur_list}
    keys.sort(
        key=lambda k: (cur_map.get(k) or {}).get("prodDataMgmt", {}).get("count", 0),
        reverse=True,
    )
    for key in keys[:16]:
        cur = cur_map.get(key) or {
            "processName": key,
            "prodDataMgmt": {"count": 0, "quantity": 0},
            "auto": {"count": 0},
            "prodDataMgmtCountRatio": 0,
        }
        cmp = cmp_map.get(key) or {"prodDataMgmt": {"count": 0}}
        c_cnt = int(cur.get("prodDataMgmt", {}).get("count", 0) or 0)
        p_cnt = int(cmp.get("prodDataMgmt", {}).get("count", 0) or 0)
        diff = c_cnt - p_cnt
        body.append(
            [
                str(cur.get("processName") or key),
                _fmt_num(c_cnt),
                _fmt_num(p_cnt),
                f"{'+' if diff > 0 else ''}{_fmt_num(diff)}",
                _fmt_qty(cur.get("prodDataMgmt", {}).get("quantity")),
                _fmt_pct(cur.get("prodDataMgmtCountRatio")),
                _fmt_num(cur.get("auto", {}).get("count")),
            ]
        )

    n = max(len(body) - 1, 1)
    _add_table(
        slide,
        body,
        Inches(0.35),
        Inches(0.8),
        Inches(12.6),
        Inches(min(0.48 + 0.36 * n, 6.3)),
        header_color=PRIMARY,
        body_size=14,
        col_widths=[1.8, 1.5, 1.5, 1.2, 1.8, 1.3, 1.5],
    )


def _build_trend_table_slide(prs: Presentation, data: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ROSE)
    trend = data.get("byMonthTrend") or []
    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.2),
        Inches(12.5),
        Inches(0.45),
        "月次推移一覧",
        font_size=28,
        bold=True,
        color=PRIMARY,
    )
    header = [
        "月",
        "修正件数",
        "集計件数",
        "総件数",
        "修正比率",
        "修正数量",
        "数量比率",
    ]
    body: List[List[str]] = [header]
    for r in trend:
        body.append(
            [
                str(r.get("month") or ""),
                _fmt_num(r.get("prodDataMgmtCount")),
                _fmt_num(r.get("autoCount")),
                _fmt_num(r.get("totalCount")),
                _fmt_pct(r.get("prodDataMgmtCountRatio")),
                _fmt_qty(r.get("prodDataMgmtQuantity")),
                _fmt_pct(r.get("prodDataMgmtQuantityRatio")),
            ]
        )
    n = max(len(body) - 1, 1)
    _add_table(
        slide,
        body,
        Inches(0.5),
        Inches(0.9),
        Inches(12.3),
        Inches(min(0.5 + 0.42 * n, 6.2)),
        header_color=PRIMARY,
        body_size=16,
        col_widths=[1.4, 1.5, 1.5, 1.4, 1.5, 2.0, 1.5],
    )


def build_manual_entry_statistics_pptx(data: Dict[str, Any]) -> bytes:
    """実績修正統計 PPT バイナリを生成する。"""
    prs = _new_prs()
    _build_cover(prs, data)
    _build_kpi_slide(prs, data)
    _build_month_compare_slide(prs, data)
    _build_trend_slide(prs, data)
    _build_process_chart_slide(prs, data)
    _build_process_table_slide(prs, data)
    _build_trend_table_slide(prs, data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
