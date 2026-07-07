"""Extract text content from PPTX for analysis."""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_shape_text(shape, depth=0):
    lines = []
    indent = "  " * depth
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            lines.append(f"{indent}{text}")
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            lines.extend(extract_shape_text(child, depth + 1))
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(f"{indent}| {' | '.join(cells)} |")
    return lines


def main(pptx_path: str, out_path: str):
    prs = Presentation(pptx_path)
    buf = []
    buf.append(f"Total slides: {len(prs.slides)}")
    buf.append(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
    buf.append("=" * 80)
    for i, slide in enumerate(prs.slides, 1):
        buf.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            lines = extract_shape_text(shape)
            for line in lines:
                buf.append(line)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                buf.append(f"  [Notes] {notes}")
    Path(out_path).write_text("\n".join(buf), encoding="utf-8")
    print(f"Written to {out_path}")


if __name__ == "__main__":
    base = Path(__file__).resolve().parents[1]
    pptx = sys.argv[1] if len(sys.argv) > 1 else str(base / "7月生産検討会資料.pptx")
    out = sys.argv[2] if len(sys.argv) > 2 else str(base / "scripts" / "pptx_content.txt")
    main(pptx, out)
