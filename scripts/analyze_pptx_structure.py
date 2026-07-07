"""Analyze PPTX structure for beautification."""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


def analyze(pptx_path: str):
    prs = Presentation(pptx_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n=== Slide {i} (layout: {slide.slide_layout.name}) ===")
        for j, shape in enumerate(slide.shapes):
            info = f"  [{j}] {shape.shape_type} name={shape.name!r}"
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().replace("\n", " | ")[:80]
                info += f" text={text!r}"
            if hasattr(shape, "width"):
                info += f" pos=({shape.left.inches:.2f},{shape.top.inches:.2f}) size=({shape.width.inches:.2f}x{shape.height.inches:.2f})"
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                info += f" table={shape.table.rows.__len__()}x{len(shape.table.columns)}"
            lines.append(info)
    Path("scripts/pptx_structure.txt").write_text("\n".join(lines), encoding="utf-8")
    print("done")


if __name__ == "__main__":
    base = Path(__file__).resolve().parents[1]
    analyze(str(base / "7月生産検討会資料.pptx"))
