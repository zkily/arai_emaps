"""生産検討会資料 PPT 生成。

画面（ProductionReviewManagement）に表示されている payload データから
ワイド画面スライドを新規生成する。テンプレート差し込みは行わない。
"""
from __future__ import annotations

import base64
import io
import re
from typing import Any, Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 画面に近いコーポレート配色
PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_LIGHT = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2D, 0x34, 0x40)
TEXT_MUTED = RGBColor(0x5A, 0x6A, 0x7A)
TABLE_ALT = RGBColor(0xF4, 0xF7, 0xFB)
TABLE_BORDER = RGBColor(0xC5, 0xD3, 0xE0)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
ORANGE = RGBColor(0xC2, 0x6A, 0x00)
SCRAP_ORANGE = RGBColor(0xF9, 0x73, 0x16)
SCRAP_VIOLET = RGBColor(0x7C, 0x3A, 0xED)
SCRAP_ROSE = RGBColor(0xE1, 0x1D, 0x48)
SCRAP_GREEN = RGBColor(0x16, 0xA3, 0x4A)
SCRAP_RED = RGBColor(0xDC, 0x26, 0x26)
SCRAP_CARD_BG = RGBColor(0xFF, 0xF7, 0xED)
FONT = "Meiryo UI"


def _fmt_th(value: Any, digits: int = 1) -> str:
    try:
        return f"{float(value):,.{digits}f}"
    except (TypeError, ValueError):
        return "—"


def _fmt_delta(value: Any, digits: int = 1) -> str:
    try:
        v = float(value)
    except (TypeError, ValueError):
        return "—"
    if v > 0:
        return f"{v:.{digits}f}"
    if v < 0:
        return f"△{abs(v):.{digits}f}"
    return f"{v:.{digits}f}"


def _fmt_prod_delta(value: Any) -> str:
    try:
        v = int(round(float(value)))
    except (TypeError, ValueError):
        return "—"
    if v > 0:
        return f"+{v}"
    if v < 0:
        return f"△{abs(v)}"
    return "+0"


def _fmt_days(value: Any) -> str:
    try:
        return f"{float(value):.1f}"
    except (TypeError, ValueError):
        return "—"


def _fmt_pct(value: Any, digits: int = 0) -> str:
    try:
        return f"{float(value):.{digits}f}%"
    except (TypeError, ValueError):
        return "—"


def _clean_comments(comments: Any) -> List[str]:
    if not comments:
        return []
    out: List[str] = []
    for c in comments:
        s = str(c or "").strip()
        if s:
            out.append(s)
    return out


def _iter_inv_rows_flat(rows: Sequence[Dict[str, Any]], depth: int = 0):
    for row in rows or []:
        yield depth, row
        for child in row.get("children") or []:
            yield depth + 1, child


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def _set_run(run, *, size: int, bold: bool = False, color: RGBColor = TEXT_DARK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=font_size, bold=bold, color=color)
    return box


def _add_accent_bar(slide, color: RGBColor = ACCENT) -> None:
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        Inches(0.12),
        SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _style_table(table, header_color: RGBColor = PRIMARY) -> None:
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    if r_idx == 0:
                        run.font.size = Pt(10)
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.size = Pt(9)
                        run.font.bold = False
                        run.font.color.rgb = TEXT_DARK
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                fill.fore_color.rgb = TABLE_ALT
            else:
                fill.fore_color.rgb = WHITE


def _add_table(
    slide,
    rows: List[List[str]],
    left,
    top,
    width,
    height,
    *,
    header_color: RGBColor = PRIMARY,
    col_widths: Optional[Sequence[float]] = None,
):
    if not rows or not rows[0]:
        return None
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    if col_widths and len(col_widths) == len(rows[0]):
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(int(width * (w / total)))
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            # 工程名は左寄せ
            if c_idx == 0:
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
    _style_table(table, header_color=header_color)
    # 左寄せを再適用（style で中央にした後）
    for r_idx, row in enumerate(rows):
        cell = table.cell(r_idx, 0)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = FONT
                if r_idx == 0:
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.size = Pt(9)
                    run.font.color.rgb = TEXT_DARK
    return table


def _add_comments(
    slide,
    comments: Sequence[str],
    *,
    left,
    top,
    width,
    max_lines: int = 6,
    title: str = "コメント",
) -> None:
    lines = _clean_comments(comments)[:max_lines]
    _add_textbox(slide, left, top, width, Inches(0.28), title, font_size=12, bold=True, color=PRIMARY)
    y = top + Inches(0.32)
    if not lines:
        _add_textbox(
            slide,
            left,
            y,
            width,
            Inches(0.3),
            "（コメント未入力）",
            font_size=10,
            color=TEXT_MUTED,
        )
        return
    for c in lines:
        _add_textbox(slide, left, y, width, Inches(0.32), f"■ {c}", font_size=10, color=TEXT_DARK)
        y += Inches(0.34)


def _build_cover(prs: Presentation, meta: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    # 背景帯
    banner = slide.shapes.add_shape(1, Inches(0), Inches(2.2), SLIDE_W, Inches(2.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY
    banner.line.fill.background()

    meeting_date = meta.get("meeting_date", "")
    if meeting_date:
        try:
            y, m, d = str(meeting_date).split("-")
            date_label = f"{y}年{int(m)}月{int(d)}日"
        except ValueError:
            date_label = str(meeting_date)
    else:
        date_label = ""

    _add_textbox(
        slide,
        Inches(1),
        Inches(2.5),
        Inches(11.3),
        Inches(0.9),
        meta.get("title") or "生産検討会",
        font_size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.5),
        Inches(3.45),
        Inches(10.3),
        Inches(0.5),
        meta.get("subtitle") or "",
        font_size=14,
        color=ACCENT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    footer = "　".join(x for x in [date_label, "生産管理部"] if x)
    _add_textbox(
        slide,
        Inches(1),
        Inches(5.3),
        Inches(11.3),
        Inches(0.4),
        footer,
        font_size=14,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    if meta.get("title_note"):
        _add_textbox(
            slide,
            Inches(1),
            Inches(5.8),
            Inches(11.3),
            Inches(0.4),
            str(meta.get("title_note")),
            font_size=11,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )


def _build_section_divider(prs: Presentation, title: str, subtitle: str, color: RGBColor = ACCENT) -> None:
    slide = _blank_slide(prs)
    bar = slide.shapes.add_shape(1, Inches(3.5), Inches(3.15), Inches(6.3), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    _add_textbox(
        slide,
        Inches(1),
        Inches(2.4),
        Inches(11.3),
        Inches(0.7),
        title or "",
        font_size=28,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.5),
        Inches(3.4),
        Inches(10.3),
        Inches(0.6),
        subtitle or "",
        font_size=14,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def _build_performance_slide(prs: Presentation, performance: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ACCENT)
    month = performance.get("month_label") or ""
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.28),
        Inches(12),
        Inches(0.45),
        f"{month} 工程別実績一覧",
        font_size=22,
        bold=True,
        color=PRIMARY,
    )
    header = [
        "工程名",
        "工程計画",
        "実見",
        "実績",
        "対実見",
        "対計画",
        "前月能率",
        "当月能率",
        "増減",
    ]
    body: List[List[str]] = [header]
    for row in performance.get("rows") or []:
        key = row.get("key")
        is_ship = key == "shipping"
        body.append(
            [
                str(row.get("name") or ""),
                _fmt_th(row.get("plan_th")),
                _fmt_th(row.get("forecast_th")),
                _fmt_th(row.get("actual_th")),
                _fmt_delta(row.get("vs_forecast_th")),
                _fmt_delta(row.get("vs_plan_th")),
                "—"
                if is_ship
                else (
                    str(int(round(float(row["productivity_prev"]))))
                    if row.get("productivity_prev") is not None
                    else "—"
                ),
                "—"
                if is_ship
                else (
                    str(int(round(float(row["productivity_curr"]))))
                    if row.get("productivity_curr") is not None
                    else "—"
                ),
                "—"
                if is_ship or row.get("productivity_delta") is None
                else _fmt_prod_delta(row.get("productivity_delta")),
            ]
        )
    n = max(len(body) - 1, 1)
    _add_table(
        slide,
        body,
        Inches(0.4),
        Inches(0.85),
        Inches(12.5),
        Inches(min(0.38 + 0.32 * n, 4.2)),
        col_widths=[1.4, 1.1, 1.1, 1.1, 1.1, 1.1, 1.3, 1.3, 1.0],
    )
    _add_comments(
        slide,
        performance.get("comments") or [],
        left=Inches(0.5),
        top=Inches(5.25),
        width=Inches(12.3),
        title="実績コメント",
    )


def _build_scrap_slide(prs: Presentation, scrap: Dict[str, Any]) -> None:
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ORANGE)
    title = "廃棄率及び廃棄本数"
    fiscal = scrap.get("fiscal_year_label") or ""
    range_label = scrap.get("range_label") or ""
    sub = " ".join(x for x in [fiscal, range_label] if x)
    _add_textbox(slide, Inches(0.5), Inches(0.28), Inches(10), Inches(0.4), title, font_size=22, bold=True, color=PRIMARY)
    if sub:
        _add_textbox(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.28), sub, font_size=11, color=TEXT_MUTED)

    monthly = scrap.get("monthly") or []
    if monthly:
        labels = []
        for m in monthly:
            yy = m.get("year")
            mm = m.get("month")
            if yy and mm:
                labels.append(f"{int(yy) % 100}/{int(mm)}")
            else:
                labels.append(f"{m.get('month', '')}月")
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series(
            "廃棄率（新）(%)",
            [float(m.get("rate_new_pct") or 0) for m in monthly],
        )
        chart_data.add_series(
            "廃棄率（旧）(%)",
            [float(m.get("rate_old_pct", m.get("rate_pct") or 0)) for m in monthly],
        )
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.35),
            Inches(1.05),
            Inches(6.2),
            Inches(3.3),
            chart_data,
        )
        chart_data2 = CategoryChartData()
        chart_data2.categories = labels
        chart_data2.add_series(
            "廃棄本数(千本)",
            [float(m.get("loss_th", m.get("scrap_th") or 0)) for m in monthly],
        )
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(6.7),
            Inches(1.05),
            Inches(6.1),
            Inches(3.3),
            chart_data2,
        )

    def _f(key: str, digits: int = 2) -> str:
        try:
            return f"{float(scrap.get(key) or 0):.{digits}f}"
        except (TypeError, ValueError):
            return "0"

    summary = (
        f"月平均廃棄率（新） {_f('avg_rate_new_current_fy_pct')}%  /  "
        f"前年度 {_f('avg_rate_new_prev_fy_pct')}%\n"
        f"月平均廃棄率（旧） {_f('avg_rate_old_current_fy_pct')}%  /  "
        f"前年度 {_f('avg_rate_old_prev_fy_pct')}%\n"
        f"当月廃棄本数 {int(scrap.get('current_month_loss_qty') or 0):,} 本"
    )
    _add_textbox(slide, Inches(0.5), Inches(4.5), Inches(6.2), Inches(1.0), summary, font_size=11, color=TEXT_DARK)
    _add_comments(
        slide,
        scrap.get("comments") or [],
        left=Inches(6.9),
        top=Inches(4.5),
        width=Inches(5.8),
        max_lines=5,
        title="廃棄コメント",
    )


def _inventory_title(inv: Dict[str, Any], *, forecast: bool) -> str:
    label = inv.get("inventory_month_label") or ""
    if forecast:
        return f"{label} 在庫予測"
    return f"{label} 月末在庫"


def _build_inventory_slide(
    prs: Presentation,
    inv: Dict[str, Any],
    *,
    forecast: bool = False,
    header_color: RGBColor = TEAL,
) -> None:
    if not inv:
        return
    slide = _blank_slide(prs)
    _add_accent_bar(slide, header_color)
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.25),
        Inches(12),
        Inches(0.4),
        _inventory_title(inv, forecast=forecast),
        font_size=22,
        bold=True,
        color=PRIMARY,
    )
    header_line = (
        f"{inv.get('prev_forecast_label', '')}出荷内示 {_fmt_th(inv.get('prev_forecast_th'))} 千本  |  "
        f"{inv.get('curr_forecast_label', '')}出荷内示 {_fmt_th(inv.get('curr_forecast_th'))} 千本"
    )
    _add_textbox(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.28), header_line, font_size=11, color=TEXT_MUTED)

    if forecast:
        header = ["工程名", "前月在庫", "前月日数", "予測在庫", "当月日数", "増減"]
    else:
        curr_hdr = "当月在庫"
        as_of = inv.get("curr_inventory_as_of")
        if as_of and isinstance(as_of, str) and len(as_of) >= 10:
            try:
                curr_hdr = f"当月在庫 {int(as_of[5:7])}/{int(as_of[8:10])}"
            except ValueError:
                pass
        header = ["工程名", "前月在庫", "前月日数", curr_hdr, "当月日数", "増減"]

    body: List[List[str]] = [header]
    for depth, row in _iter_inv_rows_flat(inv.get("rows") or []):
        indent = "　" * depth
        name = f"{indent}{row.get('name') or ''}"
        curr_val = row.get("curr_inventory_th")
        body.append(
            [
                name,
                _fmt_th(row.get("prev_inventory_th")),
                _fmt_days(row.get("prev_days")),
                _fmt_th(curr_val),
                _fmt_days(row.get("curr_days")),
                _fmt_delta(row.get("delta_th")),
            ]
        )

    n = max(len(body) - 1, 1)
    _add_table(
        slide,
        body,
        Inches(0.4),
        Inches(1.05),
        Inches(12.5),
        Inches(min(0.36 + 0.28 * n, 4.0)),
        header_color=header_color,
        col_widths=[2.2, 1.6, 1.4, 1.8, 1.4, 1.4],
    )

    # 製品 KPI（月末在庫のみ）
    if not forecast:
        product = next((r for r in (inv.get("rows") or []) if r.get("key") == "product"), None)
        if product:
            level = inv.get("product_level") or ""
            kpi = (
                f"製品 補正在庫率 {_fmt_th(product.get('curr_rate_adj'), 2)} "
                f"(目標 {_fmt_th(inv.get('product_target_rate'), 2)})  /  "
                f"在庫日数 {_fmt_days(product.get('curr_days'))}日 "
                f"(目標 {_fmt_days(inv.get('product_target_days'))}日)"
                + (f"  [{level}]" if level else "")
            )
            _add_textbox(
                slide,
                Inches(0.5),
                Inches(5.15),
                Inches(12.3),
                Inches(0.28),
                kpi,
                font_size=11,
                color=TEXT_DARK,
            )

    comment_top = Inches(5.45) if not forecast else Inches(5.2)
    _add_comments(
        slide,
        inv.get("comments") or [],
        left=Inches(0.5),
        top=comment_top,
        width=Inches(12.3),
        max_lines=4,
        title="在庫予測コメント" if forecast else "在庫コメント",
    )


def _build_load_plan_slide(prs: Presentation, load_plan: Dict[str, Any]) -> None:
    if not load_plan:
        return
    slide = _blank_slide(prs)
    _add_accent_bar(slide, ACCENT)
    month = load_plan.get("month_label") or ""
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(0.25),
        Inches(12),
        Inches(0.4),
        f"{month} 生産計画・負荷率",
        font_size=22,
        bold=True,
        color=PRIMARY,
    )
    wd = load_plan.get("working_days", 0)
    header_line = (
        f"稼働日：{wd}日  |  "
        f"内示：{_fmt_th(load_plan.get('forecast_th'))} 千本"
        f"（日当 {_fmt_th(load_plan.get('daily_forecast_th'))} 千本）"
    )
    _add_textbox(slide, Inches(0.5), Inches(0.68), Inches(12), Inches(0.28), header_line, font_size=11, color=TEXT_MUTED)

    header = [
        "工程",
        "計画",
        "日当",
        "直",
        "負荷率",
        "設備稼働率",
        "日均H",
        "設備",
        "能率",
        "稼働日",
        "定時H",
        "所要H",
    ]
    body: List[List[str]] = [header]
    for row in load_plan.get("rows") or []:
        util = row.get("equipment_utilization_pct")
        util_txt = "—" if util is None or row.get("process_cd") == "inspection" else _fmt_pct(util)
        body.append(
            [
                str(row.get("process_name") or ""),
                _fmt_th(row.get("plan_th")),
                _fmt_th(row.get("daily_th")),
                str(row.get("shift_label") or "—"),
                _fmt_pct(row.get("load_rate_pct")),
                util_txt,
                str(row.get("daily_operation_hours") if row.get("daily_operation_hours") is not None else "—"),
                str(row.get("equipment_label") or "—"),
                str(row.get("standard_rate") if row.get("standard_rate") is not None else "—"),
                str(row.get("working_days") if row.get("working_days") is not None else "—"),
                str(row.get("regular_hours") if row.get("regular_hours") is not None else "—"),
                str(row.get("required_hours") if row.get("required_hours") is not None else "—"),
            ]
        )
    n = max(len(body) - 1, 1)
    _add_table(
        slide,
        body,
        Inches(0.3),
        Inches(1.0),
        Inches(12.7),
        Inches(min(0.36 + 0.3 * n, 3.8)),
        col_widths=[1.2, 1.0, 0.9, 0.7, 0.9, 1.1, 0.8, 1.2, 0.8, 0.8, 0.8, 0.8],
    )
    _add_comments(
        slide,
        load_plan.get("comments") or [],
        left=Inches(0.5),
        top=Inches(5.15),
        width=Inches(12.3),
        max_lines=5,
        title="計画コメント",
    )


def build_production_review_pptx(data: Dict[str, Any]) -> bytes:
    """画面 payload と同じ構造の data から PPTX を生成する。"""
    prs = _new_prs()
    meta = data.get("meta") or {}
    part01 = data.get("part01") or {}
    part02 = data.get("part02") or {}
    part03 = data.get("part03") or {}

    _build_cover(prs, meta)

    _build_section_divider(
        prs,
        part01.get("title") or "PART 01",
        part01.get("subtitle") or "",
        ACCENT,
    )
    _build_performance_slide(prs, part01.get("performance") or {})
    _build_scrap_slide(prs, part01.get("scrap") or {})
    _build_inventory_slide(prs, part01.get("inventory") or {}, forecast=False, header_color=TEAL)

    _build_section_divider(
        prs,
        part02.get("title") or "PART 02",
        part02.get("subtitle") or "",
        ACCENT,
    )
    _build_load_plan_slide(prs, part02.get("load_plan") or {})
    if part02.get("inventory_forecast"):
        _build_inventory_slide(
            prs,
            part02.get("inventory_forecast") or {},
            forecast=True,
            header_color=TEAL,
        )

    _build_section_divider(
        prs,
        part03.get("title") or "PART 03",
        part03.get("subtitle") or "",
        ACCENT,
    )
    _build_load_plan_slide(prs, part03.get("load_plan") or {})
    if part03.get("inventory_forecast"):
        _build_inventory_slide(
            prs,
            part03.get("inventory_forecast") or {},
            forecast=True,
            header_color=TEAL,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _decode_chart_image(data_url: Optional[str]) -> Optional[bytes]:
    if not data_url:
        return None
    raw = str(data_url).strip()
    m = re.match(r"^data:image/(?:png|jpeg|jpg);base64,(.+)$", raw, re.I | re.S)
    payload = m.group(1) if m else raw
    try:
        return base64.b64decode(payload)
    except Exception:
        return None


def _add_rounded_rect(slide, left, top, width, height, fill: RGBColor, *, line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.18
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def _fill_shape_text(
    shape,
    lines: Sequence[Tuple[str, int, bool, RGBColor]],
    *,
    align=PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    if not lines:
        tf.text = ""
        return
    tf.text = lines[0][0]
    p0 = tf.paragraphs[0]
    p0.alignment = align
    p0.space_before = Pt(0)
    p0.space_after = Pt(2)
    for run in p0.runs:
        _set_run(run, size=lines[0][1], bold=lines[0][2], color=lines[0][3])
    for text, size, bold, color in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        _set_run(run, size=size, bold=bold, color=color)


def _scrap_loss_qty(m: Dict[str, Any]) -> int:
    if m.get("loss_qty") is not None:
        try:
            return int(m.get("loss_qty") or 0)
        except (TypeError, ValueError):
            pass
    for key in ("loss_th", "scrap_th"):
        if m.get(key) is not None:
            try:
                return int(round(float(m.get(key) or 0) * 1000))
            except (TypeError, ValueError):
                pass
    return 0


def _scrap_rate(m: Dict[str, Any], kind: str) -> float:
    key = "rate_new_pct" if kind == "new" else "rate_old_pct"
    try:
        return float(m.get(key, m.get("rate_pct") or 0) or 0)
    except (TypeError, ValueError):
        return 0.0


def _build_scrap_dashboard_slide(
    prs: Presentation,
    scrap: Dict[str, Any],
    *,
    chart_png: Optional[bytes] = None,
    meeting_label: str = "",
) -> None:
    """画面「廃棄率及び廃棄本数」と同レイアウトの 1 スライド。"""
    slide = _blank_slide(prs)

    # 背景カード風（薄いオレンジ）
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.25), Inches(0.18), Inches(12.83), Inches(7.15))
    bg.adjustments[0] = 0.04
    bg.fill.solid()
    bg.fill.fore_color.rgb = SCRAP_CARD_BG
    bg.line.color.rgb = RGBColor(0xFD, 0xBA, 0x74)
    bg.line.width = Pt(1.5)

    # 上部アクセント
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.18), Inches(12.83), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = SCRAP_ORANGE
    top_bar.line.fill.background()

    fiscal = str(scrap.get("fiscal_year_label") or "")
    range_label = str(scrap.get("range_label") or "")
    sub_parts = [x for x in [fiscal, range_label, meeting_label] if x]
    _add_textbox(
        slide,
        Inches(0.45),
        Inches(0.32),
        Inches(10),
        Inches(0.4),
        "廃棄率及び廃棄本数",
        font_size=24,
        bold=True,
        color=PRIMARY,
    )
    if sub_parts:
        _add_textbox(
            slide,
            Inches(0.45),
            Inches(0.72),
            Inches(10),
            Inches(0.28),
            "  ·  ".join(sub_parts),
            font_size=12,
            color=TEXT_MUTED,
        )

    # KPI 5 枚（画面と同じ並び）
    monthly = list(scrap.get("monthly") or [])
    last = monthly[-1] if monthly else {}
    current_new = scrap.get("current_month_rate_new_pct")
    if current_new is None:
        current_new = _scrap_rate(last, "new")
    current_old = scrap.get("current_month_rate_old_pct")
    if current_old is None:
        current_old = _scrap_rate(last, "old")
    current_loss = scrap.get("current_month_loss_qty")
    if current_loss is None:
        current_loss = _scrap_loss_qty(last)
    avg_new = scrap.get("avg_rate_new_current_fy_pct")
    avg_old = scrap.get("avg_rate_old_current_fy_pct")
    avg_loss = scrap.get("avg_loss_current_fy_qty")
    if avg_new is None and monthly:
        avg_new = sum(_scrap_rate(m, "new") for m in monthly) / len(monthly)
    if avg_old is None and monthly:
        avg_old = sum(_scrap_rate(m, "old") for m in monthly) / len(monthly)
    if avg_loss is None and monthly:
        avg_loss = sum(_scrap_loss_qty(m) for m in monthly) / len(monthly)
    prev_new = float(scrap.get("avg_rate_new_prev_fy_pct") or 0)
    prev_old = float(scrap.get("avg_rate_old_prev_fy_pct") or scrap.get("avg_rate_prev_fy_pct") or 0)
    try:
        improve_new = prev_new - float(avg_new or 0)
    except (TypeError, ValueError):
        improve_new = 0.0
    try:
        improve_old = prev_old - float(avg_old or 0)
    except (TypeError, ValueError):
        improve_old = 0.0

    kpis = [
        (
            SCRAP_VIOLET,
            "当月廃棄率（新）",
            f"{float(current_new or 0):.2f}%",
            f"年度平均 {float(avg_new or 0):.2f}%",
        ),
        (
            SCRAP_ROSE,
            "当月廃棄率（旧）",
            f"{float(current_old or 0):.2f}%",
            f"年度平均 {float(avg_old or 0):.2f}%",
        ),
        (
            SCRAP_ORANGE,
            "当月廃棄本数",
            f"{int(current_loss or 0):,}本",
            f"年度平均 {int(round(float(avg_loss or 0))):,} 本",
        ),
        (
            SCRAP_GREEN if improve_new >= 0 else SCRAP_RED,
            "廃棄率（新）改善 pt",
            f"{improve_new:.2f}pt",
            "前年度比（低いほど良）",
        ),
        (
            SCRAP_GREEN if improve_old >= 0 else SCRAP_RED,
            "廃棄率（旧）改善 pt",
            f"{improve_old:.2f}pt",
            "前年度比（低いほど良）",
        ),
    ]
    kpi_left = 0.45
    kpi_gap = 0.18
    kpi_w = (12.4 - kpi_gap * 4) / 5
    for i, (color, label, value, hint) in enumerate(kpis):
        left = Inches(kpi_left + i * (kpi_w + kpi_gap))
        card = _add_rounded_rect(slide, left, Inches(1.05), Inches(kpi_w), Inches(0.95), color)
        _fill_shape_text(
            card,
            [
                (label, 9, True, WHITE),
                (value, 18, True, WHITE),
                (hint, 8, False, RGBColor(0xFF, 0xF7, 0xED)),
            ],
            align=PP_ALIGN.LEFT,
        )

    # 計算式説明
    _add_textbox(
        slide,
        Inches(0.45),
        Inches(2.1),
        Inches(12.4),
        Inches(0.28),
        "新：主ライン（切断〜検査）連乗ロス率  |  旧：全工程（不良＋廃棄）÷ 切断実績  |  本：不良＋廃棄本数",
        font_size=10,
        color=TEXT_MUTED,
    )

    # チャート（画面 ECharts 画像があればそれを優先）
    chart_top = Inches(2.4)
    chart_h = Inches(2.85)
    if chart_png:
        stream = io.BytesIO(chart_png)
        slide.shapes.add_picture(stream, Inches(0.45), chart_top, width=Inches(12.4), height=chart_h)
    elif monthly:
        labels = [f"{int(m.get('month') or 0)}月" for m in monthly]
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("廃棄本数", [_scrap_loss_qty(m) for m in monthly])
        chart_data.add_series("廃棄率（新）(%)", [_scrap_rate(m, "new") for m in monthly])
        chart_data.add_series("廃棄率（旧）(%)", [_scrap_rate(m, "old") for m in monthly])
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.45),
            chart_top,
            Inches(12.4),
            chart_h,
            chart_data,
        )
    else:
        _add_textbox(
            slide,
            Inches(0.45),
            chart_top + Inches(1.0),
            Inches(12.4),
            Inches(0.4),
            "月度データがありません",
            font_size=14,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )

    # 月次カード
    if monthly:
        n = len(monthly)
        gap = 0.08
        card_w = min(1.05, (12.4 - gap * max(n - 1, 0)) / max(n, 1))
        total_w = n * card_w + (n - 1) * gap
        start_x = 0.45 + max(0, (12.4 - total_w) / 2)
        y = 5.4
        for i, m in enumerate(monthly):
            is_curr = i == n - 1
            left = Inches(start_x + i * (card_w + gap))
            fill = RGBColor(0xFF, 0xED, 0xD5) if is_curr else WHITE
            line = SCRAP_ORANGE if is_curr else RGBColor(0xE2, 0xE8, 0xF0)
            card = _add_rounded_rect(slide, left, Inches(y), Inches(card_w), Inches(0.92), fill, line=line)
            _fill_shape_text(
                card,
                [
                    (f"{int(m.get('month') or 0)}月", 9, True, TEXT_MUTED),
                    (f"{_scrap_rate(m, 'new'):.2f}%", 10, True, SCRAP_VIOLET),
                    (f"{_scrap_rate(m, 'old'):.2f}%", 9, True, SCRAP_ROSE),
                    (f"{_scrap_loss_qty(m):,}本", 10, True, SCRAP_ORANGE),
                ],
                align=PP_ALIGN.CENTER,
            )

    # コメント
    comments = _clean_comments(scrap.get("comments") or [])
    _add_textbox(
        slide,
        Inches(0.45),
        Inches(6.45),
        Inches(12.4),
        Inches(0.25),
        "廃棄コメント",
        font_size=12,
        bold=True,
        color=PRIMARY,
    )
    if comments:
        text = "\n".join(f"■ {c}" for c in comments[:4])
        _add_textbox(slide, Inches(0.45), Inches(6.7), Inches(12.4), Inches(0.55), text, font_size=10, color=TEXT_DARK)
    else:
        _add_textbox(
            slide,
            Inches(0.45),
            Inches(6.7),
            Inches(12.4),
            Inches(0.3),
            "（コメント未入力）",
            font_size=10,
            color=TEXT_MUTED,
        )


def build_scrap_section_pptx(
    scrap: Dict[str, Any],
    *,
    chart_image_base64: Optional[str] = None,
    meeting_label: str = "",
) -> bytes:
    """廃棄率及び廃棄本数セクション専用 PPT（画面レイアウト準拠）。"""
    prs = _new_prs()
    chart_png = _decode_chart_image(chart_image_base64)
    _build_scrap_dashboard_slide(
        prs,
        scrap or {},
        chart_png=chart_png,
        meeting_label=meeting_label,
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
